# Copyright (c) 2025 Cherry Studio MCP Tools Contributors
# SPDX-License-Identifier: CC-BY-NC-4.0
#
# 本项目采用 CC BY-NC 4.0 许可协议 — 仅限非商业用途
# 严禁任何形式的商业使用。完整许可文本见项目根目录 LICENSE 文件。
# Commercial use is STRICTLY PROHIBITED. See LICENSE for details.
"""
Cherry Office MCP Server — Word/Excel/PPT 读写
依赖: pip install mcp python-docx openpyxl python-pptx
"""
import asyncio
from pathlib import Path
from mcp.server import Server
from mcp.server.stdio import stdio_server
from mcp.types import Tool, TextContent

server = Server("cherry-office")

@server.list_tools()
async def list_tools():
    return [
        Tool(name="word_read", description="读取 .docx 文件，返回结构化文本（含段落和表格）。",
             inputSchema={"type":"object","properties":{"file_path":{"type":"string","description":"Word文件路径"}},"required":["file_path"]}),
        Tool(name="word_write", description="写入 .docx 文件，支持标题和正文。",
             inputSchema={"type":"object","properties":{"output_path":{"type":"string"},"title":{"type":"string"},"sections":{"type":"array","items":{"type":"object","properties":{"heading":{"type":"string"},"body":{"type":"string"}}}}},"required":["output_path","sections"]}),
        Tool(name="excel_read", description="读取 .xlsx/.xls 文件，返回指定 Sheet 数据。",
             inputSchema={"type":"object","properties":{"file_path":{"type":"string"},"sheet":{"type":"string","description":"Sheet名称，默认第一个"},"max_rows":{"type":"integer","description":"最大行数，默认200"}},"required":["file_path"]}),
        Tool(name="excel_write", description="写入 .xlsx 文件，支持多 Sheet。",
             inputSchema={"type":"object","properties":{"output_path":{"type":"string"},"sheets":{"type":"array","items":{"type":"object","properties":{"name":{"type":"string"},"headers":{"type":"array","items":{"type":"string"}},"rows":{"type":"array","items":{"type":"array","items":{"type":"string"}}}}}}},"required":["output_path","sheets"]}),
        Tool(name="ppt_read", description="读取 .pptx 文件，提取每页幻灯片文本。",
             inputSchema={"type":"object","properties":{"file_path":{"type":"string"}},"required":["file_path"]}),
        Tool(name="ppt_write", description="写入 .pptx 文件，每页一个标题+内容。",
             inputSchema={"type":"object","properties":{"output_path":{"type":"string"},"slides":{"type":"array","items":{"type":"object","properties":{"title":{"type":"string"},"content":{"type":"array","items":{"type":"string"}}}}}},"required":["output_path","slides"]}),
    ]

@server.call_tool()
async def call_tool(name: str, arguments: dict):
    try:
        fn = {
            "word_read": _word_read, "word_write": _word_write,
            "excel_read": _excel_read, "excel_write": _excel_write,
            "ppt_read": _ppt_read, "ppt_write": _ppt_write,
        }.get(name)
        if fn:
            return await fn(arguments)
        return [TextContent(type="text", text=f"❌ 未知工具: {name}")]
    except Exception as e:
        return [TextContent(type="text", text=f"❌ {name} 错误: {str(e)}")]

async def _word_read(args):
    from docx import Document
    path = Path(args["file_path"])
    if not path.exists():
        return [TextContent(type="text", text=f"❌ 文件不存在: {path}")]
    doc = Document(str(path))
    out = [f"📄 {path.name}\n"]
    for p in doc.paragraphs:
        if p.style.name.startswith("Heading"):
            out.append(f"\n## {p.text}")
        elif p.text.strip():
            out.append(p.text)
    for i, table in enumerate(doc.tables):
        out.append(f"\n── 表格{i+1} ──")
        for row in table.rows[:20]:
            out.append(" | ".join(cell.text[:50] for cell in row.cells))
    return [TextContent(type="text", text="\n".join(out[:5000]))]

async def _word_write(args):
    from docx import Document
    doc = Document()
    if args.get("title"):
        doc.add_heading(args["title"], 0)
    for sec in args.get("sections", []):
        doc.add_heading(sec.get("heading", ""), level=1)
        doc.add_paragraph(sec.get("body", ""))
    path = Path(args["output_path"])
    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(path))
    return [TextContent(type="text", text=f"✅ Word 已保存: {path}\n📏 {path.stat().st_size/1024:.1f} KB")]

async def _excel_read(args):
    from openpyxl import load_workbook
    path = Path(args["file_path"])
    if not path.exists():
        return [TextContent(type="text", text=f"❌ 文件不存在: {path}")]
    wb = load_workbook(str(path), read_only=True, data_only=True)
    sheet = wb[args.get("sheet", wb.sheetnames[0])]
    max_rows = int(args.get("max_rows", 200))
    out = [f"📊 {path.name} → Sheet: {sheet.title} | {sheet.max_row}行×{sheet.max_column}列\n"]
    for row in sheet.iter_rows(max_row=max_rows, values_only=True):
        out.append("\t".join(str(c) if c is not None else "" for c in row))
    wb.close()
    return [TextContent(type="text", text="\n".join(out[:5000]))]

async def _excel_write(args):
    from openpyxl import Workbook
    wb = Workbook()
    for i, sheet_def in enumerate(args.get("sheets", [])):
        ws = wb.active if i == 0 else wb.create_sheet(title=sheet_def.get("name", f"Sheet{i+1}"))
        if sheet_def.get("headers"):
            ws.append(sheet_def["headers"])
        for row in sheet_def.get("rows", []):
            ws.append(row)
    path = Path(args["output_path"])
    path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(path))
    return [TextContent(type="text", text=f"✅ Excel 已保存: {path}\n📏 {path.stat().st_size/1024:.1f} KB")]

async def _ppt_read(args):
    from pptx import Presentation
    path = Path(args["file_path"])
    if not path.exists():
        return [TextContent(type="text", text=f"❌ 文件不存在: {path}")]
    prs = Presentation(str(path))
    out = [f"📽️ {path.name} | {len(prs.slides)} 页幻灯片\n"]
    for i, slide in enumerate(prs.slides):
        out.append(f"\n── 第{i+1}页 ──")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        out.append(para.text)
    return [TextContent(type="text", text="\n".join(out[:5000]))]

async def _ppt_write(args):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    for slide_def in args.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_def.get("title", "")
        body = slide.placeholders[1]
        body.text = "\n".join(slide_def.get("content", []))
    path = Path(args["output_path"])
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return [TextContent(type="text", text=f"✅ PPT 已保存: {path}\n📏 {path.stat().st_size/1024:.1f} KB")]

async def main():
    async with stdio_server() as (read_stream, write_stream):
        await server.run(read_stream, write_stream, server.create_initialization_options())

if __name__ == "__main__":
    asyncio.run(main())
